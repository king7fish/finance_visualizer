# Finance Dashboard Ali v11 — Comparative Intelligence Edition
# Streamlit single-file app (free to host on Streamlit Community Cloud)
# Features:
# - Multi-file + multi-sheet compare
# - In-sheet multi-table detection (split a sheet into Table 1 / 2 / 3)
# - Optional AI column matching (fuzzy header mapping)
# - Optional Real-time chart suggestion
# - Axis labels, prefixes/suffixes, k/M/B scales, color pickers
# - Smart date alignment (D/W/M/Q/Y), normalization, adaptive downsample
# - Insights (MoM/YoY + correlation + quick compare)
# - Crash-proof exports (CSV/XLSX/JSON + PNG/PDF/PPTX with fallback)
# - Snapshot link

import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import re, json, base64, warnings
from io import BytesIO
from dateutil import parser
from pptx import Presentation
from pptx.util import Inches
from docx import Document
from PyPDF2 import PdfReader
from PIL import Image, ImageDraw

# ------------------ PAGE CONFIG & STYLE ------------------
st.set_page_config(page_title="Finance Dashboard Ali v11", layout="wide")
warnings.filterwarnings("ignore", category=UserWarning)

st.markdown("""
<style>
.block-container { max-width: 1240px; }
.card {
  border: 1px solid #e5e7eb; border-radius: 12px; padding: 18px 20px;
  background: #ffffff; box-shadow: 0 1px 2px rgba(0,0,0,0.05); margin-bottom: 18px;
}
h2,h3,h4 { color:#111827; margin-top:4px; }
.stButton button { border-radius:8px !important; height:42px !important; font-size:15px !important; }
.small { color:#6b7280; font-size:13px; }
</style>
""", unsafe_allow_html=True)

st.markdown(
    "<h1 style='text-align:center;color:#2563EB;margin-bottom:0;'>Finance Dashboard Ali v11</h1>"
    "<p style='text-align:center;color:#6b7280;font-size:17px;'>Multi-source comparison & in-sheet intelligence — free to use on the web.</p>",
    unsafe_allow_html=True,
)

# ------------------ SESSION FLAGS ------------------
if "disable_image_exports" not in st.session_state:
    st.session_state["disable_image_exports"] = False

# ------------------ UTILITIES ------------------
def _strip(s):
    return str(s).strip().replace("\n", " ").replace("\xa0", " ").replace("\ufeff", "")

@st.cache_data(show_spinner=False)
def safe_to_datetime(series: pd.Series) -> pd.Series:
    """Fast datetime parser with fallback to dateutil."""
    try:
        return pd.to_datetime(series, errors="coerce", infer_datetime_format=True)
    except Exception:
        out = []
        for v in series:
            try:
                out.append(parser.parse(str(v)))
            except Exception:
                out.append(pd.NaT)
        return pd.Series(out)

def looks_numeric(s: pd.Series) -> bool:
    vals = s.dropna().astype(str).head(100)
    patt = re.compile(r"^\s*[-+]?\d{1,3}(?:,\d{3})*(?:\.\d+)?\s*$|^\s*[-+]?\d+(?:\.\d+)?\s*$")
    return sum(bool(patt.match(v)) for v in vals) >= 0.6 * len(vals) if len(vals) else False

def looks_date(s: pd.Series) -> bool:
    vals = s.dropna().astype(str).head(100)
    if len(vals) == 0: return False
    # presence of separators or month names is a strong hint
    mnths = ("jan","feb","mar","apr","may","jun","jul","aug","sep","oct","nov","dec")
    score = 0
    for v in vals:
        vlow = v.lower()
        if "/" in v or "-" in v or any(m in vlow for m in mnths):
            score += 1
    return score >= 0.5 * len(vals)

@st.cache_data(show_spinner=False)
def smart_clean_dataframe(df_in: pd.DataFrame) -> pd.DataFrame:
    """Trim headers, coerce numeric/date, strip text, drop empty rows."""
    df = df_in.dropna(how="all").copy()
    df.columns = [_strip(c) for c in df.columns]
    for col in df.columns:
        s = df[col].astype(str)
        if looks_numeric(s):
            s = s.str.replace(",", "", regex=False)
            s = s.str.replace(r"[^0-9.\-]", "", regex=True).replace("", np.nan)
            df[col] = pd.to_numeric(s, errors="coerce")
        elif looks_date(s):
            df[col] = safe_to_datetime(s)
        else:
            df[col] = s.str.strip()
    return df.reset_index(drop=True)

def detect_types(df: pd.DataFrame):
    date_cols, num_cols, txt_cols = [], [], []
    for c in df.columns:
        s = df[c]
        if pd.api.types.is_datetime64_any_dtype(s) or looks_date(s.astype(str)): date_cols.append(c)
        elif pd.api.types.is_numeric_dtype(s) or looks_numeric(s.astype(str)):   num_cols.append(c)
        else:                                                                    txt_cols.append(c)
    return date_cols, num_cols, txt_cols

# ------------------ IN-SHEET MULTI-TABLE DETECTION ------------------
def split_sheet_into_tables(df: pd.DataFrame, min_blank_run: int = 1):
    """
    Split a single sheet DataFrame into multiple tables by blocks of blank rows
    or by detecting header-like rows (many non-null strings).
    Returns list of (table_df, start_row, end_row).
    """
    # If highly regular, return as one table
    if df.dropna(how="all").empty:
        return []

    # Identify blank rows
    blank = df.isna().all(axis=1)
    # Identify header-like rows: many non-null values and many strings
    header_like = []
    for i, row in df.iterrows():
        nonnull = row.notna().sum()
        strs = sum(isinstance(x, str) for x in row.values)
        header_like.append(nonnull >= max(2, int(df.shape[1]*0.4)) and strs >= max(1, int(df.shape[1]*0.3)))
    header_like = pd.Series(header_like, index=df.index)

    # Mark split indices where blank run >= min_blank_run
    blocks = []
    start = 0
    i = 0
    n = len(df)
    while i < n:
        # find next blank run
        if blank.iloc[i]:
            run = 1
            j = i+1
            while j < n and blank.iloc[j]:
                run += 1; j += 1
            if run >= min_blank_run:
                # end block at i-1
                if start <= i-1:
                    blocks.append((start, i-1))
                # next block starts at j
                start = j
                i = j
                continue
        i += 1
    if start <= n-1:
        blocks.append((start, n-1))

    # Try to refine each block by header-like shift: if first row is header-like -> use as header
    tables = []
    for (a, b) in blocks:
        seg = df.iloc[a:b+1].copy()
        seg = seg.reset_index(drop=True)
        # try to promote first non-empty row as header if looks like header
        head_idx = None
        for r in range(min(3, len(seg))):
            row = seg.iloc[r]
            nonnull = row.notna().sum()
            strs = sum(isinstance(x, str) for x in row.values)
            if nonnull >= max(2, int(df.shape[1]*0.4)) and strs >= max(1, int(df.shape[1]*0.3)):
                head_idx = r; break
        if head_idx is not None:
            new_cols = [ _strip(x) if pd.notna(x) else "" for x in list(seg.iloc[head_idx].values) ]
            body = seg.iloc[head_idx+1:].copy()
            body.columns = new_cols
            tables.append((body, a, b))
        else:
            # Use default numeric headers
            tables.append((seg, a, b))
    # Clean each table
    cleaned = []
    for t, a, b in tables:
        if not t.empty:
            t2 = smart_clean_dataframe(t)
            if not t2.dropna(how="all").empty:
                cleaned.append((t2, a, b))
    return cleaned

# ------------------ FILE UPLOAD & LOADING ------------------
with st.sidebar:
    st.header("Upload Your Data Files")
    uploaded_files = st.file_uploader(
        "Excel / CSV / JSON / DOCX / PPTX / PDF",
        type=["xlsx","xls","csv","json","docx","pptx","pdf"],
        accept_multiple_files=True
    )
    st.caption("Tip: You can upload multiple Excel files and compare sheets or even tables within a sheet.")

@st.cache_data(show_spinner=True)
def load_files(uploaded_list):
    """
    Returns a dict: key -> DataFrame
    Keys include:
      - "file.xlsx - SheetName" (for whole sheet)
      - "file.xlsx - SheetName - Table N" (for detected in-sheet tables)
      - raw tables from DOCX/PPTX and PDF text pages
    """
    tables = {}
    if not uploaded_list: return tables
    for uploaded in uploaded_list:
        name = uploaded.name
        lower = name.lower()
        try:
            if lower.endswith(("xlsx","xls")):
                xls = pd.ExcelFile(uploaded)
                for s in xls.sheet_names:
                    df_sheet = xls.parse(s, header=None)  # header=None to allow detection
                    # Split into sub tables
                    subtables = split_sheet_into_tables(df_sheet)
                    if subtables:
                        # also keep a "whole sheet" version, cleaned with first row as header if meaningful
                        df_try = xls.parse(s)
                        tables[f"{name} - {s}"] = smart_clean_dataframe(df_try)
                        for i, (tbl, a, b) in enumerate(subtables, start=1):
                            tables[f"{name} - {s} - Table {i}"] = tbl
                    else:
                        # no subtables found, parse normal
                        df_normal = xls.parse(s)
                        tables[f"{name} - {s}"] = smart_clean_dataframe(df_normal)

            elif lower.endswith("csv"):
                tables[name] = smart_clean_dataframe(pd.read_csv(uploaded))

            elif lower.endswith("json"):
                tables[name] = smart_clean_dataframe(pd.read_json(uploaded))

            elif lower.endswith("docx"):
                doc = Document(uploaded)
                for i, t in enumerate(doc.tables):
                    rows = [[cell.text for cell in row.cells] for row in t.rows]
                    df = pd.DataFrame(rows[1:], columns=rows[0]) if len(rows)>1 else pd.DataFrame(rows)
                    tables[f"{name} - Table {i+1}"] = smart_clean_dataframe(df)

            elif lower.endswith("pptx"):
                prs = Presentation(uploaded)
                for i, slide in enumerate(prs.slides):
                    for shape in slide.shapes:
                        if hasattr(shape, "table"):
                            tbl = shape.table
                            rows = [[tbl.cell(r,c).text for c in range(len(tbl.columns))] for r in range(len(tbl.rows))]
                            df = pd.DataFrame(rows[1:], columns=rows[0]) if len(rows)>1 else pd.DataFrame(rows)
                            tables[f"{name} - Slide {i+1} Table"] = smart_clean_dataframe(df)

            elif lower.endswith("pdf"):
                pdf = PdfReader(uploaded)
                pages = [pg.extract_text() for pg in pdf.pages if pg.extract_text()]
                tables[f"{name} - PDF Text"] = pd.DataFrame({"Text": pages})

        except Exception as e:
            st.error(f"File load failed for {name}: {e}")
    return tables

tables = load_files(uploaded_files) if uploaded_files else {}

# ------------------ ANALYST ENGINE: FREQ & ALIGN ------------------
def infer_pandas_freq(dates: pd.Series):
    d = safe_to_datetime(dates.dropna())
    if d.empty or d.nunique() < 3: return "M"
    try:
        inferred = pd.infer_freq(d.sort_values().unique())
    except Exception:
        inferred = None
    if not inferred:
        diffs = d.sort_values().diff().dropna()
        if diffs.empty: return "M"
        mean_days = diffs.dt.days.mean()
        if mean_days <= 2: return "D"
        if mean_days <= 9: return "W"
        if mean_days <= 40: return "M"
        if mean_days <= 100: return "Q"
        return "Y"
    inferred = inferred.upper()
    if inferred.startswith("D"): return "D"
    if inferred.startswith("W"): return "W"
    if inferred.startswith("M"): return "M"
    if inferred.startswith("Q"): return "Q"
    if inferred.startswith("A") or inferred.startswith("Y"): return "Y"
    return "M"

def to_pandas_rule(freq_code):
    return {"D":"D", "W":"W-MON", "M":"MS", "Q":"QS", "Y":"YS"}.get(freq_code, "MS")

@st.cache_data(show_spinner=True)
def resample_to_freq(df, x_col, y_cols, target_freq_code):
    out = df.copy()
    out.columns = [str(c).strip().replace("\xa0", " ").replace("\ufeff", "") for c in out.columns]
    y_existing = [y for y in y_cols if y in out.columns]
    if not y_existing or x_col not in out.columns:
        return pd.DataFrame(columns=[x_col] + list(y_existing))
    out[x_col] = safe_to_datetime(out[x_col])
    out = out.dropna(subset=[x_col])
    if out.empty:
        return pd.DataFrame(columns=[x_col] + list(y_existing))
    out = out.set_index(x_col)
    rule = to_pandas_rule(target_freq_code)
    try:
        res = out[y_existing].resample(rule).mean()
    except Exception:
        return pd.DataFrame(columns=[x_col] + list(y_existing))
    res = res.reset_index()
    res.columns = [x_col] + list(y_existing)
    return res

def align_on_calendar(dfs_info, target_freq):
    aligned = []
    master_x = "Master_X"
    for info in dfs_info:
        name = info["name"]; df = info["df"].copy()
        x_col = info["x_col"]; y_cols = info["y_cols"]; suffix = info.get("suffix","")
        df.columns = [str(c).strip().replace("\xa0", " ").replace("\ufeff", "") for c in df.columns]
        valid_y = [y for y in y_cols if y in df.columns]
        if not valid_y or x_col not in df.columns:
            st.warning(f"[{name}] skipped — missing valid X or Y columns after cleaning.")
            continue
        rs = resample_to_freq(df, x_col, valid_y, target_freq)
        if rs.empty:
            st.warning(f"[{name}] produced no rows after resampling.")
            continue
        melt = rs.melt(id_vars=[x_col], var_name="Series", value_name="Value")
        if suffix:
            melt["Series"] = melt["Series"].astype(str) + " " + suffix
        melt.rename(columns={x_col: master_x}, inplace=True)
        aligned.append(melt)
    if not aligned:
        return pd.DataFrame(columns=["Master_X", "Series", "Value"])
    return pd.concat(aligned, ignore_index=True)

def apply_normalization(df_long, mode, group_col="Series"):
    out = df_long.copy()
    out["Value"] = pd.to_numeric(out["Value"], errors="coerce")
    if mode == "Off": return out
    if mode == "Z-score":
        def zfun(s):
            v = s["Value"].astype(float)
            mu = v.mean(skipna=True); sd = v.std(skipna=True)
            s["Value"] = (v - mu) / (sd if sd and sd != 0 else 1.0)
            return s
        return out.groupby(group_col, group_keys=False).apply(zfun)
    if mode.startswith("Relative"):
        def rfun(s):
            v = s["Value"].astype(float)
            first = v.dropna().iloc[0] if v.dropna().size else np.nan
            s["Value"] = np.nan if pd.isna(first) or first == 0 else (v / first) * 100.0
            return s
        return out.groupby(group_col, group_keys=False).apply(rfun)
    return out

def adaptive_downsample(df_long, x_col="Master_X", max_points=20000):
    if df_long.shape[0] <= max_points: return df_long
    out = []
    budget = max_points // max(1, df_long["Series"].nunique())
    for s, sdf in df_long.groupby("Series"):
        sdf = sdf.sort_values(x_col)
        if sdf.shape[0] <= budget: out.append(sdf); continue
        idx = np.linspace(0, len(sdf)-1, num=budget).astype(int)
        out.append(sdf.iloc[idx])
    return pd.concat(out, ignore_index=True)

# ------------------ AI COLUMN MATCHING (optional) ------------------
AI_SYNONYMS = {
    "date": ["date","period","month","year","quarter","qtr","time","week","day","dt"],
    "value": ["value","amount","revenue","sales","profit","income","cost","price","qty","quantity","metric","count","users","visits"],
    "category": ["category","name","label","segment","product","dept","region","class"]
}
def ai_guess(colnames):
    names = [str(c).lower().strip() for c in colnames]
    # heuristic score per synonym root
    scores = {k: [] for k in AI_SYNONYMS}
    for c in names:
        for root, syns in AI_SYNONYMS.items():
            score = 0
            for s in syns:
                if s in c: score += 1
            scores[root].append(score)
    # pick best matching index for each root
    picks = {}
    for root, arr in scores.items():
        if not arr: continue
        idx = int(np.argmax(arr))
        if arr[idx] > 0: picks[root] = names[idx]
    # return the original column names
    out = {}
    for root, nm in picks.items():
        for original in colnames:
            if str(original).lower().strip() == nm:
                out[root] = original
                break
    return out

# ------------------ SNAPSHOT ------------------
def encode_snapshot(state: dict) -> str:
    try:
        raw = json.dumps(state, separators=(",",":")).encode("utf-8")
        return base64.urlsafe_b64encode(raw).decode("utf-8")
    except Exception: return ""

def decode_snapshot(token: str) -> dict:
    try:
        raw = base64.urlsafe_b64decode(token.encode("utf-8"))
        return json.loads(raw.decode("utf-8"))
    except Exception: return {}

# ------------------ TABS ------------------
tab_data, tab_viz, tab_insight, tab_export, tab_settings = st.tabs(
    ["Data", "Visualize", "Insights", "Export", "Settings"]
)

# ------------------ DATA TAB ------------------
with tab_data:
    st.markdown('<div class="card">', unsafe_allow_html=True)
    if not tables:
        st.info("Upload Excel/CSV/JSON/DOCX/PPTX/PDF files in the sidebar to begin.")
    else:
        key = st.selectbox("Preview a dataset (file/sheet/table)", list(tables.keys()))
        df_prev = tables[key].copy()
        st.success(f"Loaded {key} — {df_prev.shape[0]:,} rows × {df_prev.shape[1]} columns")
        st.dataframe(df_prev.head(15), use_container_width=True)
    st.markdown('</div>', unsafe_allow_html=True)

# ------------------ VISUALIZE TAB ------------------
with tab_viz:
    st.markdown('<div class="card">', unsafe_allow_html=True)
    st.subheader("Step 1: Choose datasets (across files/sheets or in-sheet tables)")
    if not tables:
        st.info("Upload data to visualize.")
        st.markdown('</div>', unsafe_allow_html=True)
        st.stop()

    params = st.query_params
    snap_param = params.get("snap")
    snapshot_token = snap_param[0] if isinstance(snap_param, list) else snap_param
    pre_selected = []
    if snapshot_token:
        snap = decode_snapshot(snapshot_token)
        pre_selected = snap.get("sources", [])

    compare_mode = st.checkbox("Compare Mode (overlay multiple datasets)", value=True)
    if compare_mode:
        default_src = pre_selected if pre_selected else list(tables.keys())[:2]
        sources = st.multiselect("Pick datasets/sheets/tables", list(tables.keys()),
                                 default=[s for s in default_src if s in tables][:5])
    else:
        default_one = pre_selected[0] if pre_selected else (list(tables.keys())[0] if tables else None)
        sources = [st.selectbox("Choose one dataset", list(tables.keys()),
                                index=0 if not default_one else list(tables.keys()).index(default_one))]
    st.markdown('</div>', unsafe_allow_html=True)

    # Toggles for AI features
    st.markdown('<div class="card">', unsafe_allow_html=True)
    st.subheader("Step 2: Optional assists")
    col_ai1, col_ai2 = st.columns(2)
    with col_ai1:
        enable_ai_mapping = st.checkbox("AI column matching (map Date/Value/Category)", value=False)
    with col_ai2:
        enable_chart_suggest = st.checkbox("Real-time chart suggestion", value=False,
                                           help="Suggests chart type after data prep; you can still override.")
    st.markdown('</div>', unsafe_allow_html=True)

    # AXES + LABELS + SCALES
    st.markdown('<div class="card">', unsafe_allow_html=True)
    st.subheader("Step 3: Axes, labels & chart options")
    c1, c2 = st.columns(2)
    with c1:
        custom_x_label = st.text_input("X-Axis Label", value="X")
        x_prefix = st.text_input("X Prefix", value="")
        x_suffix = st.text_input("X Suffix", value="")
    with c2:
        custom_y_label = st.text_input("Y-Axis Label", value="Value")
        y_prefix = st.text_input("Y Prefix", value="")
        y_suffix = st.text_input("Y Suffix", value="")

    srow1 = st.columns(4)
    with srow1[0]:
        chart_type = st.selectbox("Chart Type", ["Auto","Line","Area","Bar","Scatter","Pie"], index=0)
    with srow1[1]:
        x_scale = st.selectbox("X Scale", ["None","Thousands (/1,000)","Millions (/1,000,000)","Billions (/1,000,000,000)"])
    with srow1[2]:
        y_scale = st.selectbox("Y Scale", ["None","Thousands (/1,000)","Millions (/1,000,000)","Billions (/1,000,000,000)"])
    with srow1[3]:
        layout_mode = st.selectbox("Layout", ["Overlay (single chart)","Side-by-side panels"], index=0)
    st.markdown('</div>', unsafe_allow_html=True)

    # PER-SOURCE MAPPING + COLORS
    st.markdown('<div class="card">', unsafe_allow_html=True)
    st.subheader("Step 4: Column mapping + colors")
    dfs_info = []
    color_map = {}

    for src in sources:
        df_src = tables[src].copy()
        df_src.columns = [_strip(c) for c in df_src.columns]
        date_cols, num_cols, txt_cols = detect_types(df_src)

        # AI mapping suggestion
        ai_map = {}
        if enable_ai_mapping:
            try:
                ai_map = ai_guess(df_src.columns)
            except Exception:
                ai_map = {}

        b1, b2 = st.columns([1,1])
        with b1:
            x_col_default = ai_map.get("date", (date_cols[0] if date_cols else df_src.columns[0]))
            x_col = st.selectbox(f"X column for [{src}]", df_src.columns,
                                 index=list(df_src.columns).index(x_col_default) if x_col_default in df_src.columns else 0,
                                 key=f"x_{src}")
        with b2:
            label_suffix = st.text_input(f"Legend suffix for [{src}] (optional)", value="", key=f"suf_{src}")

        # AI Y candidates
        y_suggest = []
        if enable_ai_mapping and "value" in ai_map:
            y_suggest = [ai_map["value"]]
        if not y_suggest:
            y_suggest = [c for c in df_src.columns if c in num_cols][:1] or ([df_src.columns[1]] if len(df_src.columns)>1 else [])
        y_cols = st.multiselect(f"Y columns for [{src}]", df_src.columns, default=y_suggest, key=f"y_{src}")

        # Color pickers per Y series
        if y_cols:
            cols = st.columns(min(3, len(y_cols)) or 1)
            palette = ["#1f77b4","#d62728","#2ca02c","#9467bd","#ff7f0e","#8c564b","#e377c2"]
            for idx, y in enumerate(y_cols):
                with cols[idx % len(cols)]:
                    color = st.color_picker(f"Color for {y} [{src}]", value=palette[idx % len(palette)], key=f"c_{src}_{y}")
                    series_name = f"{y}{(' ' + label_suffix) if label_suffix else ''}".strip()
                    color_map[series_name] = color

        dfs_info.append({"name": src, "df": df_src, "x_col": x_col, "y_cols": y_cols, "suffix": label_suffix})
    st.markdown('</div>', unsafe_allow_html=True)

    # ALIGNMENT / NORMALIZATION / ADAPTIVE
    st.markdown('<div class="card">', unsafe_allow_html=True)
    st.subheader("Step 5: Alignment & performance")
    colA, colB, colC = st.columns(3)
    with colA:
        enable_alignment = st.checkbox("Smart date alignment (common calendar)", value=True)
    with colB:
        norm_mode = st.selectbox("Normalize scales", ["Off","Z-score","Relative to first period (%)"], index=0)
    with colC:
        adaptive_mode = st.checkbox("Adaptive mode (downsample big data)", value=True)
    max_points = st.slider("Adaptive max points across all series", 5000, 100000, 20000, step=5000)
    st.markdown('</div>', unsafe_allow_html=True)

    generate = st.button("Generate Chart", type="primary", use_container_width=True)

    # SCALE MAP
    scale_map = {"None":1, "Thousands (/1,000)":1_000, "Millions (/1,000,000)":1_000_000, "Billions (/1,000,000,000)":1_000_000_000}
    master_x_name = "Master_X"

    def auto_chart_type_from(plotted_df, chosen):
        if chosen != "Auto": return chosen
        if plotted_df.empty: return "Line"
        # if X is datetime: line is default
        if pd.api.types.is_datetime64_any_dtype(plotted_df[master_x_name]): return "Line"
        # categorical small: pie; else bar
        nunique = plotted_df[master_x_name].astype(str).nunique()
        if nunique <= 8 and plotted_df["Series"].nunique() <= 6: return "Pie"
        return "Bar"

    if generate:
        with st.spinner("Preparing and aligning data..."):
            # choose common frequency
            target_freq = "M"
            if enable_alignment:
                freqs = []
                for info in dfs_info:
                    if info["x_col"] in info["df"].columns:
                        try: freqs.append(infer_pandas_freq(info["df"][info["x_col"]]))
                        except Exception: pass
                if freqs: target_freq = pd.Series(freqs).mode().iloc[0]

            # Build long data
            if enable_alignment:
                raw_long = align_on_calendar(dfs_info, target_freq)
            else:
                combined = []
                for info in dfs_info:
                    df = info["df"]; x_col = info["x_col"]; ys = info["y_cols"]; suf = info.get("suffix","")
                    if not ys or x_col not in df.columns:
                        st.warning(f"[{info['name']}] skipped — invalid X or Y selection."); continue
                    tmp = df.copy()
                    if looks_date(tmp[x_col].astype(str)) or pd.api.types.is_datetime64_any_dtype(tmp[x_col]):
                        tmp[x_col] = safe_to_datetime(tmp[x_col]); tmp = tmp.dropna(subset=[x_col])
                    try:
                        melt = tmp[[x_col] + ys].melt(id_vars=[x_col], var_name="Series", value_name="Value")
                    except Exception as e:
                        st.warning(f"[{info['name']}] melt failed: {e}"); continue
                    if suf: melt["Series"] = melt["Series"].astype(str) + " " + suf
                    melt.rename(columns={x_col: master_x_name}, inplace=True)
                    combined.append(melt)
                raw_long = pd.concat(combined, ignore_index=True) if combined else pd.DataFrame(columns=[master_x_name,"Series","Value"])

            normalized = apply_normalization(raw_long, norm_mode, group_col="Series")

            plotted = normalized.copy()
            # (Note: X scaling only applies if numeric categories; dates won't be scaled)
            if pd.api.types.is_numeric_dtype(plotted[master_x_name]):
                plotted[master_x_name] = plotted[master_x_name] / scale_map[x_scale]
            plotted["Value"] = pd.to_numeric(plotted["Value"], errors="coerce") / scale_map[y_scale]

            if adaptive_mode and not plotted.empty:
                total_rows = int(plotted.shape[0])
                if total_rows > max_points:
                    plotted = adaptive_downsample(plotted, x_col=master_x_name, max_points=max_points)

            # Store session
            st.session_state["ali_raw_long"]   = raw_long
            st.session_state["ali_plotted"]    = plotted
            st.session_state["ali_color_map"]  = color_map
            st.session_state["ali_chart_type"] = auto_chart_type_from(plotted, chart_type if not enable_chart_suggest else "Auto")
            st.session_state["ali_layout_mode"]= layout_mode
            st.session_state["ali_labels"]     = {
                "x": f"{x_prefix}{custom_x_label}{x_suffix}",
                "y": f"{y_prefix}{custom_y_label}{y_suffix}"
            }
            st.session_state["ali_sources"]    = sources
            st.session_state["ali_target_freq"]= target_freq
            st.success("Chart data prepared. Scroll down to view.")

    # RENDER
    if "ali_plotted" in st.session_state and not st.session_state["ali_plotted"].empty:
        plotted = st.session_state["ali_plotted"]
        raw_long = st.session_state["ali_raw_long"]
        color_map = st.session_state["ali_color_map"]
        chart_choice = st.session_state["ali_chart_type"]
        layout_mode = st.session_state["ali_layout_mode"]
        labels = st.session_state["ali_labels"]

        compare_basis = None
        if chart_choice in ["Bar","Pie"]:
            st.markdown('<div class="card">', unsafe_allow_html=True)
            st.subheader("Bar/Pie comparison controls")
            c1, c2 = st.columns(2)
            with c1:
                compare_basis = st.selectbox("Comparison basis",
                    ["Totals per Series", "By X category (grouped)", "Distribution by X category (percent)"])
            with c2:
                layout_mode = st.selectbox("Layout", ["Overlay (single chart)", "Side-by-side panels"],
                                           index=0 if layout_mode.startswith("Overlay") else 1)
            st.markdown('</div>', unsafe_allow_html=True)

        # Build figure
        fig = None
        try:
            if chart_choice == "Pie":
                if compare_basis == "Totals per Series":
                    agg = plotted.groupby("Series", as_index=False)["Value"].sum()
                    fig = px.pie(agg, names="Series", values="Value", color="Series", color_discrete_map=color_map)
                elif compare_basis == "By X category (grouped)":
                    cats = plotted[master_x_name].dropna().astype(str).unique().tolist()
                    if cats:
                        chosen_cat = st.selectbox("Pick X category for pie", cats[:200])
                        slice_df = plotted[plotted[master_x_name].astype(str) == chosen_cat]
                        agg = slice_df.groupby("Series", as_index=False)["Value"].sum()
                        fig = px.pie(agg, names="Series", values="Value", color="Series",
                                     title=f"Category: {chosen_cat}", color_discrete_map=color_map)
                    else:
                        st.info("No X categories available for pie.")
                else:
                    tmp = plotted.dropna(subset=["Value"]).copy()
                    denom = tmp.groupby(master_x_name)["Value"].transform("sum")
                    tmp["SharePct"] = np.where(denom > 0, (tmp["Value"] / denom) * 100.0, np.nan)
                    agg = tmp.groupby("Series", as_index=False)["SharePct"].mean(numeric_only=True)
                    agg["SharePct"] = agg["SharePct"].fillna(0)
                    fig = px.pie(agg, names="Series", values="SharePct", color="Series",
                                 title="Average share across categories", color_discrete_map=color_map)

            elif chart_choice == "Bar":
                if compare_basis == "Totals per Series":
                    agg = plotted.groupby("Series", as_index=False)["Value"].sum()
                    fig = px.bar(agg, x="Series", y="Value", color="Series", color_discrete_map=color_map)
                elif compare_basis == "By X category (grouped)":
                    grp = plotted.groupby([master_x_name, "Series"], as_index=False)["Value"].sum()
                    if layout_mode.startswith("Overlay"):
                        fig = px.bar(grp, x=master_x_name, y="Value", color="Series",
                                     barmode="group", color_discrete_map=color_map)
                    else:
                        fig = px.bar(grp, x=master_x_name, y="Value", facet_col="Series",
                                     color="Series", color_discrete_map=color_map)
                else:
                    tmp = plotted.dropna(subset=["Value"]).copy()
                    denom = tmp.groupby(master_x_name)["Value"].transform("sum")
                    tmp["SharePct"] = np.where(denom > 0, (tmp["Value"] / denom) * 100.0, np.nan)
                    tmp["SharePct"] = tmp["SharePct"].fillna(0)
                    if layout_mode.startswith("Overlay"):
                        fig = px.bar(tmp, x=master_x_name, y="SharePct", color="Series",
                                     barmode="group", color_discrete_map=color_map)
                    else:
                        fig = px.bar(tmp, x=master_x_name, y="SharePct", facet_col="Series",
                                     color="Series", color_discrete_map=color_map)

            else:
                args = dict(x=master_x_name, y="Value", color="Series", color_discrete_map=color_map)
                if chart_choice == "Line":
                    fig = px.line(plotted, markers=True, **args)
                elif chart_choice == "Area":
                    fig = px.area(plotted, **args)
                elif chart_choice == "Scatter":
                    fig = px.scatter(plotted, **args)
                else:
                    fig = px.line(plotted, markers=True, **args)

        except Exception as e:
            st.error(f"Chart failed: {e}")

        if fig is not None:
            fig.update_layout(template="plotly_white", height=650,
                              xaxis_title=labels["x"], yaxis_title=labels["y"])
            st.plotly_chart(fig, use_container_width=True)
            st.session_state["ali_fig"] = fig

            # Quick compare summary
            st.markdown('<div class="card">', unsafe_allow_html=True)
            st.subheader("Quick Compare Summary")
            safe = plotted.copy()
            safe["Value"] = pd.to_numeric(safe["Value"], errors="coerce")
            tbl = pd.DataFrame({
                "Mean": safe.groupby("Series")["Value"].mean(),
                "Median": safe.groupby("Series")["Value"].median(),
                "Last": safe.sort_values("Master_X").groupby("Series")["Value"].last()
            }).round(3).reset_index()
            st.dataframe(tbl, use_container_width=True)
            st.markdown('</div>', unsafe_allow_html=True)

            # Snapshot link
            st.markdown('<div class="card">', unsafe_allow_html=True)
            st.subheader("Share Snapshot")
            snap_state = {
                "sources": st.session_state.get("ali_sources", []),
                "compare_mode": compare_mode,
                "chart_type": chart_choice,
                "layout_mode": layout_mode,
                "enable_alignment": enable_alignment,
                "norm_mode": norm_mode,
                "adaptive_mode": adaptive_mode,
                "max_points": max_points,
                "x_label": labels["x"], "y_label": labels["y"],
                "x_scale": x_scale, "y_scale": y_scale
            }
            token = encode_snapshot(snap_state)
            try:
                base = st.get_url()
                share_url = f"{base}?snap={token}"
            except Exception:
                share_url = f"?snap={token}"
            st.text_input("Copy this link to restore this view:", value=share_url)
            st.markdown('</div>', unsafe_allow_html=True)

        else:
            st.info("Select at least one valid Y column to render a chart.")

# ------------------ INSIGHTS TAB ------------------
with tab_insight:
    st.markdown('<div class="card">', unsafe_allow_html=True)
    if "ali_plotted" not in st.session_state or st.session_state["ali_plotted"].empty:
        st.info("Generate a chart in the Visualize tab to see insights.")
    else:
        plotted = st.session_state["ali_plotted"]
        x_col, series_col, value_col = "Master_X", "Series", "Value"

        def yoy_mom_insights(df_long):
            bullets = []
            if df_long.empty or not pd.api.types.is_datetime64_any_dtype(df_long[x_col]):
                return ["No date-based insights available."]
            for sname, sdf in df_long.groupby(series_col):
                sdf = sdf.dropna(subset=[value_col]).sort_values(x_col)
                if sdf.empty: bullets.append(f"{sname}: no numeric values after cleaning."); continue
                monthly = sdf.set_index(x_col)[value_col].resample("M").mean()
                if len(monthly) >= 2:
                    mom = (monthly.iloc[-1] - monthly.iloc[-2]) / (abs(monthly.iloc[-2]) + 1e-9) * 100
                    bullets.append(f"{sname}: last month change {mom:+.1f}%")
                if len(monthly) >= 13 and monthly.iloc[-13] != 0:
                    yoy = (monthly.iloc[-1] - monthly.iloc[-13]) / (abs(monthly.iloc[-13]) + 1e-9) * 100
                    bullets.append(f"{sname}: year over year change {yoy:+.1f}%")
            return bullets or ["Not enough data for MoM/YoY."]

        def comparative_intelligence(df_long):
            if df_long[series_col].nunique() < 2:
                return ["Single series; no cross-series comparison."], None
            wide = df_long.pivot_table(index=x_col, columns=series_col, values=value_col, aggfunc="mean").sort_index()
            corr = None
            try: corr = wide.corr().round(3)
            except Exception: corr = None
            bullets = []
            names = list(wide.columns)
            if len(names) == 2:
                a, b = names
                a_vals, b_vals = wide[a], wide[b]
                diff = (a_vals - b_vals); pct = (a_vals - b_vals) / (b_vals.abs() + 1e-9) * 100
                avg_diff = diff.mean(skipna=True); avg_pct = pct.mean(skipna=True)
                bullets.append(f"{a} vs {b}: avg gap {avg_diff:,.2f}; avg pct gap {avg_pct:+.1f}%")
            else:
                bullets.append("Multiple series; see correlation table for relationships.")
            return bullets, corr

        st.subheader("Narrative Insights")
        kt = yoy_mom_insights(plotted)
        st.markdown("<br>".join(f"- {x}" for x in kt), unsafe_allow_html=True)

        st.subheader("Comparative Summary")
        bullets, corr = comparative_intelligence(plotted)
        st.markdown("<br>".join(f"- {x}" for x in bullets), unsafe_allow_html=True)

        st.subheader("Correlation (numeric series)")
        if corr is not None and not corr.empty:
            st.dataframe(corr, use_container_width=True)
        else:
            st.info("Not enough numeric overlap for correlation matrix.")
    st.markdown('</div>', unsafe_allow_html=True)

# ------------------ EXPORT TAB ------------------
def placeholder_png(text="No Chart Available", color=(0,0,0)):
    img = Image.new("RGB", (900, 560), color=(255, 255, 255))
    d = ImageDraw.Draw(img); d.text((280, 260), text, fill=color)
    buf = BytesIO(); img.save(buf, format="PNG"); return buf.getvalue()

def fig_to_image_safe(fig, fmt="png"):
    if fig is None: return placeholder_png("No Chart Available", (0,0,0))
    if st.session_state.get("disable_image_exports", False):
        return placeholder_png("Image export disabled", (255,165,0))
    try:
        return fig.to_image(format=fmt, engine="kaleido")
    except Exception as e:
        st.warning(f"{fmt.upper()} export fallback used: {e}")
        return placeholder_png(f"{fmt.upper()} export unavailable", (255,0,0))

def pptx_with_chart_failsafe(fig, title="Chart"):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1)).text = title
    png = fig_to_image_safe(fig, "png") if fig is not None else None
    if png:
        slide.shapes.add_picture(BytesIO(png), Inches(0.8), Inches(1.4), width=Inches(8.8))
    else:
        slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(9), Inches(3)).text_frame.text = "No chart available."
    out = BytesIO(); prs.save(out); out.seek(0); return out.getvalue()

with tab_export:
    st.markdown('<div class="card">', unsafe_allow_html=True)
    if "ali_plotted" not in st.session_state or st.session_state["ali_plotted"].empty:
        st.info("Generate a chart first.")
    else:
        fig_to_export = st.session_state.get("ali_fig")
        df_plot_long = st.session_state["ali_plotted"]
        df_raw_long = st.session_state["ali_raw_long"]

        st.subheader("Export Data")
        which = st.radio("Choose data to export:", ["Plotted (filtered + scaled)", "Raw (aligned/cleaned)"], horizontal=True)
        export_df = df_plot_long if which.startswith("Plotted") else df_raw_long
        export_df = export_df if isinstance(export_df, pd.DataFrame) else pd.DataFrame()

        c1, c2, c3 = st.columns(3)
        with c1:
            st.download_button("Download CSV", export_df.to_csv(index=False).encode("utf-8"),
                               "data.csv", "text/csv")
        with c2:
            buf = BytesIO()
            with pd.ExcelWriter(buf, engine="openpyxl") as w:
                export_df.to_excel(w, index=False, sheet_name="Data")
            st.download_button("Download Excel", buf.getvalue(),
                               "data.xlsx",
                               "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
        with c3:
            st.download_button("Download JSON", export_df.to_json(orient="records").encode("utf-8"),
                               "data.json", "application/json")

        st.markdown("---")
        st.subheader("Export Chart (crash-proof)")
        png_bytes  = fig_to_image_safe(fig_to_export, "png")
        pdf_bytes  = fig_to_image_safe(fig_to_export, "pdf")
        pptx_bytes = pptx_with_chart_failsafe(fig_to_export, title="Finance Dashboard - Chart")

        d1, d2, d3 = st.columns(3)
        with d1: st.download_button("Download PNG", png_bytes, "chart.png", "image/png")
        with d2: st.download_button("Download PDF", pdf_bytes, "chart.pdf", "application/pdf")
        with d3: st.download_button("Download PPTX", pptx_bytes, "chart_slide.pptx",
                                    "application/vnd.openxmlformats-officedocument.presentationml.presentation")
    st.markdown('</div>', unsafe_allow_html=True)

# ------------------ SETTINGS TAB ------------------
with tab_settings:
    st.markdown('<div class="card">', unsafe_allow_html=True)
    st.checkbox("Disable image exports (if Kaleido/Chrome not available)",
                value=st.session_state["disable_image_exports"],
                key="disable_image_exports")
    st.markdown("""
**Notes**
- Compare Mode overlays multiple files, sheets, or in-sheet tables (Table 1, Table 2, …).
- Smart Alignment resamples to a common calendar (D/W/M/Q/Y).
- Normalization: Off, Z-score, Relative-to-first-period (%).
- Adaptive Mode down-samples big data without losing shape.
- Full exports: CSV, Excel, JSON, PNG, PDF, PPTX (with fallbacks).
- Snapshot link restores view & settings.
""")
    st.markdown('</div>', unsafe_allow_html=True)
